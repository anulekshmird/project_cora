"""
ocr_engine.py — CORA OCR Engine
================================================
Provides text extraction via two strategies:

1. NATIVE extraction  — for Word / PDF / PPTX files where the file
   path is known. Always preferred over screenshot OCR.

2. SCREENSHOT OCR     — Tesseract with a preprocessing pipeline
   tuned for each document type (dense text, code, general UI).

Usage:
    from ocr_engine import extract_text, extract_from_file, OCRMode

    # Screenshot path (auto-selects best mode)
    text = extract_text(pil_image, mode=OCRMode.DOCUMENT)

    # File path (native — no Tesseract needed)
    text = extract_from_file("/path/to/report.docx")
"""

import re
import io
import os
from enum import Enum
from PIL import Image, ImageFilter, ImageEnhance, ImageOps


# ── OCR backend (Tesseract via pytesseract) ──────────────────────────────────
try:
    import pytesseract
    # Path to Tesseract executable discovered on this machine
    pytesseract.pytesseract.tesseract_cmd = r'C:\Users\ADITHYA\AppData\Local\Programs\Tesseract-OCR\tesseract.exe'
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    print("OCR Engine: pytesseract not found. Screenshot OCR disabled.")

# ── Native document parsers ───────────────────────────────────────────────────
try:
    import pypdf
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

try:
    import docx as python_docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ─────────────────────────────────────────────────────────────────────────────
# OCR Modes
# ─────────────────────────────────────────────────────────────────────────────

class OCRMode(Enum):
    """
    Controls the preprocessing pipeline and Tesseract config.

    AUTO        — Detect best mode from image content (default)
    DOCUMENT    — Dense body text: Word, PDF, articles
    CODE        — Source code, terminals, monospace text
    GENERAL     — Mixed UI, browser, general desktop
    SUBTITLE    — Video subtitles, short bright-on-dark text
    """
    AUTO     = "auto"
    DOCUMENT = "document"
    CODE     = "code"
    GENERAL  = "general"
    SUBTITLE = "subtitle"


# Tesseract PSM modes:
#   3  = fully automatic page segmentation (default)
#   4  = single column of text
#   6  = single uniform block of text
#   11 = sparse text (find as much text as possible)
_TESSERACT_CONFIGS = {
    OCRMode.DOCUMENT: "--psm 4 --oem 3",
    OCRMode.CODE:     "--psm 6 --oem 3 -c preserve_interword_spaces=1",
    OCRMode.GENERAL:  "--psm 3 --oem 3",
    OCRMode.SUBTITLE: "--psm 11 --oem 3",
    OCRMode.AUTO:     "--psm 3 --oem 3",
}


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def extract_text(image: Image.Image, mode: OCRMode = OCRMode.AUTO) -> str:
    """
    Extract text from a PIL screenshot image.

    Parameters
    ----------
    image : PIL.Image
        The screenshot to process.
    mode  : OCRMode
        Preprocessing + Tesseract tuning profile.
        AUTO will pick the best profile based on image content.

    Returns
    -------
    str — extracted text, cleaned of noise.
    """
    if not TESSERACT_AVAILABLE:
        return ""
    if image is None:
        return ""

    try:
        if mode == OCRMode.AUTO:
            mode = _detect_mode(image)

        processed = _preprocess(image, mode)
        config_str = _TESSERACT_CONFIGS[mode]
        raw = pytesseract.image_to_string(processed, config=config_str)
        return _clean(raw)

    except Exception as e:
        print(f"OCR extract_text error: {e}")
        return ""


def extract_from_file(path: str) -> str:
    """
    Extract text natively from a document file.
    Preferred over screenshot OCR — perfect fidelity, no image noise.

    Supports: .docx  .pdf  .pptx  .txt  .md  .py  .json  (and other text)

    Returns
    -------
    str — full document text, or empty string on failure.
    """
    if not path or not os.path.exists(path):
        return ""

    ext = os.path.splitext(path)[1].lower()

    try:
        if ext == ".docx":
            return _read_docx(path)
        elif ext == ".pdf":
            return _read_pdf(path)
        elif ext == ".pptx":
            return _read_pptx(path)
        else:
            # Plain text fallback
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(120_000)
    except Exception as e:
        print(f"OCR extract_from_file error ({ext}): {e}")
        return ""


def extract_text_for_window(
    image: Image.Image,
    window_title: str,
    file_path: str = "",
    mode_primary: str = "general",
) -> str:
    """
    Smart dispatcher: tries native file extraction first, falls back to OCR.

    This is the method copilot_controller / observer should call
    instead of raw extract_text().

    Parameters
    ----------
    image        : PIL screenshot (used only if native extraction fails)
    window_title : active window title (used to pick OCR mode)
    file_path    : if known, path to the open document
    mode_primary : context engine mode string

    Returns
    -------
    str — best available text for this window.
    """
    w = window_title.lower()

    # ── 1. Try native extraction if we have a file path ───────────────────
    if file_path:
        text = extract_from_file(file_path)
        if len(text.strip()) > 50:
            print(f"OCR: Native extraction → {len(text)} chars from {os.path.basename(file_path)}")
            return text[:4000]   # cap for prompt budget

    # ── 2. Choose OCR mode from window context ─────────────────────────────
    if any(k in w for k in ["word", "docs", "writer", "document", "notepad", ".docx"]):
        ocr_mode = OCRMode.DOCUMENT
    elif any(k in w for k in ["pdf", "acrobat", "foxit", "okular", "evince", ".pdf"]):
        ocr_mode = OCRMode.DOCUMENT
    elif any(k in w for k in ["code", "vscode", "pycharm", "vim", "terminal",
                               "powershell", "cmd", "bash", "jupyter"]):
        ocr_mode = OCRMode.CODE
    elif any(k in w for k in ["youtube", "netflix", "video", "vlc", "mpv"]):
        ocr_mode = OCRMode.SUBTITLE
    else:
        ocr_mode = OCRMode.GENERAL

    # ── 3. Screenshot OCR with appropriate pipeline ────────────────────────
    text = extract_text(image, mode=ocr_mode)
    print(f"OCR: Screenshot OCR ({ocr_mode.value}) → {len(text)} chars")
    return text


# ─────────────────────────────────────────────────────────────────────────────
# Preprocessing pipelines
# ─────────────────────────────────────────────────────────────────────────────

def _preprocess(image: Image.Image, mode: OCRMode) -> Image.Image:
    """Apply mode-specific image preprocessing before Tesseract."""

    # ── Normalise size ────────────────────────────────────────────────────
    # Upscale small images (Tesseract accuracy drops below ~150 DPI).
    # Never downscale below 2000px wide — keeps small fonts readable.
    img = image.copy()
    min_width = 2000
    if img.width < min_width:
        scale = min_width / img.width
        img = img.resize(
            (int(img.width * scale), int(img.height * scale)),
            Image.Resampling.LANCZOS,
        )

    if mode == OCRMode.DOCUMENT:
        return _preprocess_document(img)
    elif mode == OCRMode.CODE:
        return _preprocess_code(img)
    elif mode == OCRMode.SUBTITLE:
        return _preprocess_subtitle(img)
    else:
        return _preprocess_general(img)


def _preprocess_document(img: Image.Image) -> Image.Image:
    """
    Optimised for dense body text (Word, PDF).
    Steps: greyscale → denoise → contrast boost → sharpen → binarise
    """
    # Greyscale
    img = img.convert("L")

    # Mild denoise (smooths JPEG/screenshot compression artefacts)
    img = img.filter(ImageFilter.MedianFilter(size=3))

    # Contrast enhancement — makes light grey text on white background pop
    img = ImageEnhance.Contrast(img).enhance(2.2)

    # Sharpness — recovers edge detail lost in scaling
    img = ImageEnhance.Sharpness(img).enhance(2.0)

    # Adaptive binarisation via Otsu-like threshold
    img = _binarise(img, threshold=160)

    return img


def _preprocess_code(img: Image.Image) -> Image.Image:
    """
    Optimised for monospace code / terminal output.
    Dark backgrounds with bright text → invert before binarising.
    """
    img = img.convert("L")

    # Detect dark background (terminal)
    avg_brightness = sum(img.getdata()) / (img.width * img.height)
    if avg_brightness < 128:
        img = ImageOps.invert(img)  # dark bg → light bg for Tesseract

    img = ImageEnhance.Contrast(img).enhance(2.0)
    img = ImageEnhance.Sharpness(img).enhance(1.8)
    img = _binarise(img, threshold=140)

    return img


def _preprocess_subtitle(img: Image.Image) -> Image.Image:
    """
    Optimised for video subtitles — typically short white text on dark bg.
    Crop to bottom third where subtitles live.
    """
    # Crop bottom 30% of frame — that's where subtitles appear
    h = img.height
    img = img.crop((0, int(h * 0.70), img.width, h))

    img = img.convert("L")
    avg_brightness = sum(img.getdata()) / (img.width * img.height)
    if avg_brightness < 128:
        img = ImageOps.invert(img)

    img = ImageEnhance.Contrast(img).enhance(2.5)
    img = _binarise(img, threshold=150)

    return img


def _preprocess_general(img: Image.Image) -> Image.Image:
    """Balanced preprocessing for mixed UI / browser content."""
    img = img.convert("L")
    img = ImageEnhance.Contrast(img).enhance(1.8)
    img = ImageEnhance.Sharpness(img).enhance(1.5)
    return img


def _binarise(img: Image.Image, threshold: int = 155) -> Image.Image:
    """Convert greyscale image to pure black/white at given threshold."""
    return img.point(lambda p: 255 if p > threshold else 0, "1").convert("L")


def _detect_mode(image: Image.Image) -> OCRMode:
    """
    Heuristic: sample the image to pick the best OCR mode.
    Dark background → CODE, bright uniform → DOCUMENT, else GENERAL.
    """
    small = image.copy().convert("L")
    small.thumbnail((200, 200))
    pixels = list(small.getdata())
    avg = sum(pixels) / len(pixels)
    dark_ratio = sum(1 for p in pixels if p < 80) / len(pixels)

    if dark_ratio > 0.55:
        return OCRMode.CODE      # Dark terminal / code editor
    elif avg > 200:
        return OCRMode.DOCUMENT  # Bright document / Word / PDF
    else:
        return OCRMode.GENERAL


# ─────────────────────────────────────────────────────────────────────────────
# Native document readers
# ─────────────────────────────────────────────────────────────────────────────

def _read_docx(path: str) -> str:
    """
    Extract text from a .docx file with structure preservation.
    Includes headings, body paragraphs, and table cells.
    """
    if not DOCX_AVAILABLE:
        return "[python-docx not installed]"

    doc = python_docx.Document(path)
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Preserve heading hierarchy
        if para.style.name.startswith("Heading"):
            level = para.style.name.replace("Heading ", "")
            parts.append(f"\n{'#' * int(level) if level.isdigit() else '#'} {text}")
        else:
            parts.append(text)

    # Extract table content
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return "\n".join(parts)


def _read_pdf(path: str) -> str:
    """
    Extract text from a PDF with page markers.
    Falls back to a warning if the PDF is scanned (image-only).
    """
    if not PYPDF_AVAILABLE:
        return "[pypdf not installed]"

    reader = pypdf.PdfReader(path)
    parts = []

    for i, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        page_text = page_text.strip()
        if page_text:
            parts.append(f"[Page {i + 1}]\n{page_text}")

    full = "\n\n".join(parts)

    if len(full.strip()) < 50:
        return (
            "[PDF appears to contain scanned images only. "
            "Text extraction requires OCR on the PDF image layers.]"
        )

    return full


def _read_pptx(path: str) -> str:
    """Extract text from a .pptx with slide markers."""
    if not PPTX_AVAILABLE:
        return "[python-pptx not installed]"

    prs = Presentation(path)
    parts = []

    for i, slide in enumerate(prs.slides):
        slide_parts = [f"[Slide {i + 1}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)


# ─────────────────────────────────────────────────────────────────────────────
# Text cleanup
# ─────────────────────────────────────────────────────────────────────────────

def _clean(text: str) -> str:
    """Remove OCR noise and normalise whitespace."""
    if not text:
        return ""

    # Remove form-feed and other control characters
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)

    # Collapse 3+ blank lines to 2
    text = re.sub(r"\n{3,}", "\n\n", text)

    # Remove lines that are pure noise (single chars, lone punctuation)
    lines = [ln for ln in text.splitlines() if len(ln.strip()) > 1]
    text = "\n".join(lines)

    return text.strip()